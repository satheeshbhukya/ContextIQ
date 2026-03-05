import fitz  # PyMuPDF
import docx
import pandas as pd
from pptx import Presentation
from typing import BinaryIO


class DocumentParser:

    @staticmethod
    def extract_text_from_pdf(file: BinaryIO) -> str:
        try:
            pdf_document = fitz.open(stream=file.read(), filetype="pdf")
            text = []
            for page in pdf_document:
                text.append(page.get_text())
            return "\n".join(text)
        except Exception as e:
            raise ValueError(f"Error extracting text from PDF: {str(e)}")
    
    @staticmethod
    def extract_text_from_word(file: BinaryIO) -> str:
        try:
            doc = docx.Document(file)
            full_text = []
            
            # Extract paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    full_text.append(paragraph.text)
            
            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                    if row_text:
                        full_text.append(row_text)
            
            return "\n".join(full_text)
        except Exception as e:
            raise ValueError(f"Error extracting text from Word: {str(e)}")
    
    @staticmethod
    def extract_text_from_excel(file: BinaryIO) -> str:
        try:
            df = pd.read_excel(file)
            return df.to_string(index=False)
        except Exception as e:
            raise ValueError(f"Error extracting text from Excel: {str(e)}")
    
    @staticmethod
    def extract_text_from_ppt(file: BinaryIO) -> str:
        try:
            presentation = Presentation(file)
            text = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame"):
                        for paragraph in shape.text_frame.paragraphs:
                            text.append(paragraph.text)
            return "\n".join(text)
        except Exception as e:
            raise ValueError(f"Error extracting text from PowerPoint: {str(e)}")
    
    @staticmethod
    def extract_text_from_txt(file: BinaryIO) -> str:
        try:
            return file.read().decode('utf-8')
        except Exception as e:
            raise ValueError(f"Error extracting text from TXT: {str(e)}")
    
    @classmethod
    def parse_document(cls, file: BinaryIO, file_type: str) -> str:
        parsers = {
            'pdf': cls.extract_text_from_pdf,
            'docx': cls.extract_text_from_word,
            'xlsx': cls.extract_text_from_excel,
            'pptx': cls.extract_text_from_ppt,
            'txt': cls.extract_text_from_txt
        }
        
        parser = parsers.get(file_type.lower())
        if not parser:
            raise ValueError(f"Unsupported file type: {file_type}")
        
        return parser(file)
